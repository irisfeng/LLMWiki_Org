"""Text extraction for uploaded files.

Kept out of the upload request handler so extraction runs inside the
Celery worker instead of blocking the FastAPI event loop.
"""
from __future__ import annotations

import logging
import os

logger = logging.getLogger(__name__)


def _extract_pptx_text(file_path: str) -> str:
    """Walk pptx shapes directly — markitdown's PptxConverter crashes on
    shapes whose position is None (Emu vs NoneType)."""
    from pptx import Presentation
    prs = Presentation(file_path)
    parts: list[str] = []
    for i, slide in enumerate(prs.slides, 1):
        slide_parts: list[str] = []
        for shape in slide.shapes:
            try:
                if hasattr(shape, "text") and shape.text and shape.text.strip():
                    slide_parts.append(shape.text.strip())
            except Exception:
                continue
        if slide_parts:
            parts.append(f"## Slide {i}\n\n" + "\n\n".join(slide_parts))
    return "\n\n".join(parts)


def extract_text_sync(file_path: str, raw_bytes: bytes | None = None) -> str:
    """Synchronous extraction for non-PDF formats. For PDF use parse_pdf
    from mineru_service (async HTTP to a remote service)."""
    ext = os.path.splitext(file_path)[1].lower()
    if ext in (".pptx", ".ppt"):
        try:
            text = _extract_pptx_text(file_path)
            if text:
                return text
        except Exception as e:
            logger.error("pptx direct extraction failed for %s: %s", file_path, e)
    if raw_bytes is None:
        with open(file_path, "rb") as f:
            raw_bytes = f.read()
    try:
        return raw_bytes.decode("utf-8")
    except UnicodeDecodeError:
        pass
    try:
        from markitdown import MarkItDown
        return MarkItDown().convert(file_path).text_content or ""
    except Exception as e:
        logger.error("markitdown conversion failed for %s: %s", file_path, e)
        return ""


async def extract_text(file_path: str, raw_bytes: bytes | None = None) -> str:
    """Async entry point. Uses MinerU for PDF; falls back to extract_text_sync."""
    ext = os.path.splitext(file_path)[1].lower()
    if ext == ".pdf":
        try:
            from app.services.mineru_service import parse_pdf
            text = await parse_pdf(file_path)
            if text:
                return text
        except Exception as e:
            logger.error("MinerU PDF parse failed for %s: %s", file_path, e)
    return extract_text_sync(file_path, raw_bytes=raw_bytes)
