import os
import uuid
import logging
import ipaddress
import socket
from urllib.parse import urlparse
from fastapi import APIRouter, UploadFile, File, Form, Depends, HTTPException
from fastapi.responses import FileResponse
from sqlalchemy.ext.asyncio import AsyncSession
from sqlalchemy import select, func, update, delete

from app.database import get_db
from app.models import RawSource, WikiPage, WikiChunk, WikiLink
from app.schemas import SourceTextSubmit, SourceURLSubmit, SourceResponse, WikiPageSummary
from app.config import settings
from app.worker import process_ingest

logger = logging.getLogger(__name__)

router = APIRouter(prefix="/api/sources", tags=["sources"])


def _validate_url_safe(url: str) -> str:
    """Validate that a URL is safe to fetch (anti-SSRF).

    Only allows http/https to public IP addresses.
    Resolves DNS first to prevent DNS rebinding attacks.
    Raises HTTPException if the URL is unsafe.
    """
    parsed = urlparse(url)

    # Only allow http and https schemes
    if parsed.scheme not in ("http", "https"):
        raise HTTPException(status_code=400, detail="仅支持 http/https 链接")

    hostname = parsed.hostname
    if not hostname:
        raise HTTPException(status_code=400, detail="无效的 URL")

    # Block obvious localhost patterns
    if hostname in ("localhost", "127.0.0.1", "0.0.0.0", "::1", "[::1]"):
        raise HTTPException(status_code=400, detail="不允许访问本地地址")

    # Resolve DNS to get actual IP, then check if it's private/reserved
    try:
        resolved_ips = socket.getaddrinfo(hostname, None, socket.AF_UNSPEC, socket.SOCK_STREAM)
    except socket.gaierror:
        raise HTTPException(status_code=400, detail="无法解析域名")

    for family, _, _, _, sockaddr in resolved_ips:
        ip_str = sockaddr[0]
        try:
            ip = ipaddress.ip_address(ip_str)
        except ValueError:
            continue
        # Block any non-globally-routable address (private, loopback, link-local, reserved)
        if not ip.is_global:
            raise HTTPException(
                status_code=400,
                detail="不允许访问内网地址",
            )

    return url


def _extract_pptx_text(file_path: str) -> str:
    """Extract text from .pptx/.ppt by walking shapes directly.
    Bypasses markitdown's PptxConverter which crashes on shapes with null position (Emu vs NoneType).
    """
    from pptx import Presentation
    prs = Presentation(file_path)
    parts: list[str] = []
    for i, slide in enumerate(prs.slides, 1):
        slide_parts: list[str] = []
        for shape in slide.shapes:
            try:
                if hasattr(shape, "text") and shape.text and shape.text.strip():
                    slide_parts.append(shape.text.strip())
            except Exception:
                continue
        if slide_parts:
            parts.append(f"## Slide {i}\n\n" + "\n\n".join(slide_parts))
    return "\n\n".join(parts)


def _extract_text(file_path: str, raw_bytes: bytes | None = None) -> str:
    """Unified text extraction for upload + reingest.
    PDF -> MinerU (async caller should use parse_pdf directly)
    PPTX/PPT -> python-pptx (skip markitdown, which crashes)
    Other -> UTF-8 decode, then MarkItDown fallback.
    """
    ext = os.path.splitext(file_path)[1].lower()
    if ext in (".pptx", ".ppt"):
        try:
            text = _extract_pptx_text(file_path)
            if text:
                return text
        except Exception as e:
            logger.error("pptx direct extraction failed for %s: %s", file_path, e)
    if raw_bytes is None:
        with open(file_path, "rb") as f:
            raw_bytes = f.read()
    try:
        return raw_bytes.decode("utf-8")
    except UnicodeDecodeError:
        pass
    try:
        from markitdown import MarkItDown
        return MarkItDown().convert(file_path).text_content or ""
    except Exception as e:
        logger.error("markitdown conversion failed for %s: %s", file_path, e)
        return ""


@router.post("/upload", response_model=SourceResponse)
async def upload_file(
    file: UploadFile = File(...),
    submitted_by: str = Form(""),
    db: AsyncSession = Depends(get_db),
):
    os.makedirs(settings.raw_storage_path, exist_ok=True)
    file_id = str(uuid.uuid4())
    ext = os.path.splitext(file.filename)[1] if file.filename else ".txt"
    file_path = os.path.join(settings.raw_storage_path, f"{file_id}{ext}")

    content = await file.read()
    with open(file_path, "wb") as f:
        f.write(content)

    # Extract text content — failures here should not block the upload
    content_text = ""
    try:
        if ext.lower() == ".pdf":
            from app.services.mineru_service import parse_pdf
            content_text = await parse_pdf(file_path)
        if not content_text:
            content_text = _extract_text(file_path, raw_bytes=content)
    except Exception as e:
        logger.error("Text extraction failed for %s: %s", file.filename, e)

    source = RawSource(
        id=file_id, filename=file.filename or "upload", file_path=file_path,
        content_text=content_text, submitted_by=submitted_by,
        status="pending" if content_text else "failed",
        error_message=None if content_text else "文本提取失败，请检查文件格式",
    )
    db.add(source)
    await db.commit()
    await db.refresh(source)
    if content_text:
        process_ingest.delay(str(source.id))
    return source


@router.post("/text", response_model=SourceResponse)
async def submit_text(body: SourceTextSubmit, db: AsyncSession = Depends(get_db)):
    os.makedirs(settings.raw_storage_path, exist_ok=True)
    file_id = str(uuid.uuid4())
    file_path = os.path.join(settings.raw_storage_path, f"{file_id}.md")
    with open(file_path, "w", encoding="utf-8") as f:
        f.write(body.text)

    source = RawSource(
        id=file_id, filename=body.title or "text-input", file_path=file_path,
        content_text=body.text, submitted_by=body.submitted_by, status="pending",
    )
    db.add(source)
    await db.commit()
    await db.refresh(source)
    process_ingest.delay(str(source.id))
    return source


@router.post("/url", response_model=SourceResponse)
async def submit_url(body: SourceURLSubmit, db: AsyncSession = Depends(get_db)):
    import httpx

    # SSRF protection: validate URL before fetching
    _validate_url_safe(body.url)

    os.makedirs(settings.raw_storage_path, exist_ok=True)
    async with httpx.AsyncClient(timeout=30.0) as client:
        resp = await client.get(body.url, follow_redirects=True)
        resp.raise_for_status()

    file_id = str(uuid.uuid4())
    file_path = os.path.join(settings.raw_storage_path, f"{file_id}.html")
    with open(file_path, "wb") as f:
        f.write(resp.content)

    content_text = ""
    try:
        from markitdown import MarkItDown
        mid = MarkItDown()
        result = mid.convert(file_path)
        content_text = result.text_content
    except Exception as e:
        logger.error("URL content extraction failed for %s: %s", body.url, e)

    source = RawSource(
        id=file_id, filename=body.url, file_path=file_path,
        content_text=content_text, submitted_by=body.submitted_by,
        status="pending" if content_text else "failed",
        error_message=None if content_text else "网页内容提取失败",
    )
    db.add(source)
    await db.commit()
    await db.refresh(source)
    if content_text:
        process_ingest.delay(str(source.id))
    return source


@router.get("/", response_model=list[SourceResponse])
async def list_sources(db: AsyncSession = Depends(get_db)):
    """List all raw sources with per-source generated page counts."""
    sources_result = await db.execute(
        select(RawSource).order_by(RawSource.created_at.desc()).limit(200)
    )
    sources = list(sources_result.scalars().all())
    if not sources:
        return []

    counts_result = await db.execute(
        select(WikiPage.source_id, func.count(WikiPage.id))
        .where(WikiPage.source_id.in_([s.id for s in sources]))
        .group_by(WikiPage.source_id)
    )
    counts = {str(sid): c for sid, c in counts_result.all()}

    out = []
    for s in sources:
        resp = SourceResponse.model_validate(s, from_attributes=True)
        resp.generated_pages_count = counts.get(str(s.id), 0)
        out.append(resp)
    return out


@router.get("/{source_id}", response_model=SourceResponse)
async def get_source(source_id: str, db: AsyncSession = Depends(get_db)):
    result = await db.execute(select(RawSource).where(RawSource.id == source_id))
    source = result.scalar_one_or_none()
    if not source:
        raise HTTPException(status_code=404, detail="源不存在")
    count_result = await db.execute(
        select(func.count(WikiPage.id)).where(WikiPage.source_id == source.id)
    )
    resp = SourceResponse.model_validate(source, from_attributes=True)
    resp.generated_pages_count = count_result.scalar() or 0
    return resp


@router.get("/{source_id}/pages", response_model=list[WikiPageSummary])
async def list_source_pages(source_id: str, db: AsyncSession = Depends(get_db)):
    """List wiki pages generated from a specific raw source."""
    source = await db.get(RawSource, source_id)
    if not source:
        raise HTTPException(status_code=404, detail="源不存在")
    result = await db.execute(
        select(WikiPage).where(WikiPage.source_id == source.id).order_by(WikiPage.created_at)
    )
    return result.scalars().all()


@router.get("/{source_id}/download")
async def download_source_file(source_id: str, db: AsyncSession = Depends(get_db)):
    """Stream the original raw source file back to the user (PDF/DOCX/MD/HTML/...)."""
    source = await db.get(RawSource, source_id)
    if not source:
        raise HTTPException(status_code=404, detail="源不存在")
    if not source.file_path or not os.path.exists(source.file_path):
        raise HTTPException(status_code=404, detail="原文文件不存在")
    # Use original filename so the browser saves it sensibly
    safe_name = source.filename or os.path.basename(source.file_path)
    return FileResponse(
        path=source.file_path,
        filename=safe_name,
        media_type="application/octet-stream",
    )


@router.delete("/{source_id}")
async def delete_source(source_id: str, cascade: bool = False, db: AsyncSession = Depends(get_db)):
    """Delete a raw source and its file.
    If cascade=True, also delete generated wiki pages and their chunks.
    If cascade=False (default), generated pages are kept but unlinked.
    """
    source = await db.get(RawSource, source_id)
    if not source:
        raise HTTPException(status_code=404, detail="源不存在")

    deleted_pages_count = 0
    if cascade:
        # Find all pages generated from this source
        pages_result = await db.execute(
            select(WikiPage).where(WikiPage.source_id == source.id)
        )
        pages = pages_result.scalars().all()
        deleted_pages_count = len(pages)

        for page in pages:
            # Delete chunks for this page
            await db.execute(
                delete(WikiChunk).where(WikiChunk.page_id == page.id)
            )
            # Delete outgoing links from this page
            await db.execute(
                delete(WikiLink).where(WikiLink.from_page_id == page.id)
            )
            # Null out incoming links pointing to this page (preserve other pages' link records)
            await db.execute(
                update(WikiLink).where(WikiLink.to_page_id == page.id).values(to_page_id=None)
            )
            await db.delete(page)
    else:
        # Unlink generated pages (keep the knowledge, orphan the source_id)
        await db.execute(
            update(WikiPage).where(WikiPage.source_id == source.id).values(source_id=None)
        )

    # Remove physical file
    if source.file_path and os.path.exists(source.file_path):
        try:
            os.remove(source.file_path)
        except OSError as e:
            logger.warning("Failed to remove file %s: %s", source.file_path, e)

    await db.delete(source)
    await db.commit()
    return {
        "message": "已删除",
        "cascade": cascade,
        "deleted_pages": deleted_pages_count,
    }


@router.get("/{source_id}/preview")
async def preview_source(source_id: str, db: AsyncSession = Depends(get_db)):
    """Return truncated content_text for preview."""
    source = await db.get(RawSource, source_id)
    if not source:
        raise HTTPException(status_code=404, detail="源不存在")
    if not source.content_text:
        raise HTTPException(status_code=404, detail="该文档没有可预览的文本内容")
    # Truncate to 3000 chars for preview
    preview_text = source.content_text[:3000]
    has_more = len(source.content_text) > 3000
    return {
        "id": str(source.id),
        "filename": source.filename,
        "preview": preview_text,
        "total_length": len(source.content_text),
        "truncated": has_more,
    }


@router.get("/{source_id}/file")
async def download_file(source_id: str, db: AsyncSession = Depends(get_db)):
    source = await db.get(RawSource, source_id)
    if not source or not source.file_path:
        raise HTTPException(status_code=404, detail="文件不存在")
    if not os.path.exists(source.file_path):
        raise HTTPException(status_code=404, detail="文件已被删除")
    return FileResponse(
        source.file_path,
        filename=source.filename,
        media_type="application/octet-stream",
    )


@router.post("/{source_id}/reingest")
async def reingest_source(source_id: str, db: AsyncSession = Depends(get_db)):
    source = await db.get(RawSource, source_id)
    if not source:
        raise HTTPException(status_code=404, detail="源不存在")

    # Re-extract text from raw file if content_text is empty
    if not source.content_text and source.file_path and os.path.exists(source.file_path):
        ext = os.path.splitext(source.file_path)[1].lower()
        try:
            if ext == ".pdf":
                from app.services.mineru_service import parse_pdf
                source.content_text = await parse_pdf(source.file_path)
            if not source.content_text:
                source.content_text = _extract_text(source.file_path)
        except Exception as e:
            logger.error("Re-extraction failed for %s: %s", source.filename, e)

    if not source.content_text:
        raise HTTPException(status_code=400, detail="文本提取失败，无法重新处理")

    source.status = "pending"
    source.processed_at = None
    source.error_message = None
    await db.commit()
    process_ingest.delay(str(source.id))
    return {"message": "已重新加入处理队列"}
